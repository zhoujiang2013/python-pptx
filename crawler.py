# !/usr/bin/env python
# encoding:UTF-8
from util import request_url
import re
import os
import sys
#from __future__ import print_function
from pptx import Presentation
from pptx.util import Inches
import PIL

class Crawler(object):
	def __init__(self):
		self.main_url = "https://mp.weixin.qq.com/s?__biz=MzA3NzIwMDM3OA==&mid=209853452&idx=1&sn=bd40e9622dca2e5bd52af08bbf870861&pass_ticket=8MmcYuwV6RkFHjUHOnxmzVg%2FEhQYTM26Zg%2FO2ZpgJVGyL6ewBt5fJc%2BEsNkytOiN"
		self.media_content_pattern = re.compile('<div class="rich_media_content " id="js_content">.*?</div>',re.S)
		self.item_pattern = re.compile('<p><a href="(.*?)" target=',re.S)
		self.title_pattern = re.compile('<h2 class="rich_media_title" id="activity-name">(.*?)</h2>',re.S)
		self.elements_content_pattern = re.compile('<p style=.*?</p>',re.S)
		self.png_pattern = re.compile('data-src="(.*?)"',re.S)
		self.datatype_pattern = re.compile('data-type="(.*?)"',re.S)
		self.text_pattern = re.compile('>(.*?)</',re.S)
		self.picid_pattern = re.compile('http://mmbiz.qpic.cn/mmbiz/(.*?)/',re.S)
		self.pic_fmt_pattern = re.compile('wx_fmt=(.*?)$',re.S)
		self.data_path = "../data/"
		self.ppt_path = "../ppt/"

	def get_item_list(self):
		ret,main_page = request_url(self.main_url)
		if ret == -1 or main_page == "":
			print "Request main page failed!"
			return
		info = self.media_content_pattern.findall(main_page);
		if len(info) != 0:
			media_content = info[0]
		else:
			media_content = ""
		item_list = []
		if media_content:
			item_info = self.item_pattern.findall(main_page);
			if len(item_info) != 0:
				for item_url in item_info:
					print item_url
					item_list.append(item_url)
					title = self.get_item(item_url)
					if title== "" and item_url.find("&amp;")!=-1:
						item_url = item_url.replace("&amp;","&")
						title = self.get_item(item_url)

	def get_item(self,item_url):
		ret,item_page = request_url(item_url)
		if ret == -1 or item_page == "":
			print "Request item page failed! %s" % item_url
			return "bad"
		
		info = self.title_pattern.findall(item_page);
		if len(info) != 0:
			title = info[0].strip().replace("(ppt)","")
		else:
			title = ""
		
		item_path = self.data_path+"bak"
		ppt_file = self.ppt_path+"bak"
		if title!="":
			item_path = self.data_path + title
			ppt_file = self.ppt_path+ title + ".pptx"
		else:
			print "title is null!%s" % item_url
			return ""
		if os.path.exists(ppt_file):
			return "exist"
		info = self.media_content_pattern.findall(item_page);
		if len(info) != 0:
			media_content = info[0]
		else:
			media_content = ""
		if media_content == "":
			return
		info = self.elements_content_pattern.findall(media_content);
		element_tuple_list = []
		for element_content in info:
			if element_content.find('data-src="http://')!=-1:
				element_type = "png"
			else:
				element_type = "text"
			element_data = ""
			if element_type == "png":
				info = self.png_pattern.findall(element_content);
				if len(info) != 0:
					element_data = info[0]
					if element_data.find("wx_fmt")==-1:
						info = self.datatype_pattern.findall(element_content)
						if len(info)>0:
							element_data += "?wx_fmt=%s" % info[0]
				else:
					element_data = ""

			else:
				info = self.text_pattern.findall(element_content);
				if len(info)>3:
					element_data = ""
				else:
					element_data = "\n".join(info)
			if element_data:
				element_tuple_list.append((element_type,element_data))
		
		if len(element_tuple_list) > 0:
			if not os.path.exists(item_path):
				os.makedirs(item_path)
		text_data = ""
		picfile_list = []
		for element_tuple in element_tuple_list:
			element_type,element_data = element_tuple
			if element_type == "text":
				text_data += element_data
			else:
				if element_data:
					picfile = self.download_pic(element_data,item_path)
					if picfile!="":
						picfile_list.append(picfile)
		self.write_text_content(text_data,item_path)
		self.creat_ppt(title,ppt_file,picfile_list)
		return title
	def download_pic(self,url,path):
		ret,pic_page = request_url(url)
		if ret == -1 or pic_page=="":
			return ""
		info = self.picid_pattern.findall(url)
		if len(info)>0 and info[0]!="":
			picid = info[0]
		else:
			picid = url.replace("/","_").replace(":","_").replace(".","_").replace("?","_")
		info = self.pic_fmt_pattern.findall(url)
		if len(info)>0 and info[0]!="":
			fmt = info[0].split("&")[0]
			fmt =  fmt.split("?")[0]
		else:
			print "Get pic fmt failed!%s" % url
			return ""
		filename_bak = "../data/" + picid + "_bak.%s" % fmt
		
		fp = open(filename_bak,"w")
		fp.write(pic_page)
		fp.close()
		pil_image = PIL.Image.open(filename_bak)
		w, h = pil_image.size
		w_box = 720
		h_box = 540
		filename = path + "/" + picid + ".%s" % fmt
		self.resize(w, h, w_box, h_box, pil_image,filename)
		os.system("rm -rf %s" % filename_bak)
		return filename

	def resize(self,w, h, w_box, h_box, pil_image,outfile):
		'''
		resize a pil_image object so it will fit into
		a box of size w_box times h_box, but retain aspect ratio
		'''
		f1 = 1.0*w_box/w  # 1.0 forces float division in Python2
		f2 = 1.0*h_box/h
		factor = min([f1, f2])
		#print(f1, f2, factor)  # test
		# use best down-sizing filter
		width = int(w*factor)
		height = int(h*factor)
		out = pil_image.resize((width, height), PIL.Image.ANTIALIAS)
		out.save(outfile) 
		
	def write_text_content(self,text_data,path):
		filename = path + "/" + "text_data.txt"
		fp = open(filename,"w")
		fp.write(text_data)
		fp.close()

	def creat_ppt(self,title_content,ppt_file,picfile_list):
		prs = Presentation("default.pptx")
		title_slide_layout = prs.slide_layouts[0]
		slide = prs.slides.add_slide(title_slide_layout)
		title = slide.shapes.title
		title.text = title_content
		graph_slide_layout = prs.slide_layouts[6]
		for picfile in picfile_list:
			slide = prs.slides.add_slide(graph_slide_layout)
			slide.shapes.add_picture(picfile,0,0)
		prs.save(ppt_file)

def test_ppt():
	picfile = "test.jpeg"
	for i in xrange(12):
		try:
			prs = Presentation("default.pptx")
			graph_slide_layout = prs.slide_layouts[i]
			slide = prs.slides.add_slide(graph_slide_layout)
			placeholder = slide.placeholders[0]
			#pic = placeholder.insert_picture(picfile)
			prs.save("../ppt/%s.pptx" % i)
		except:
			continue

def test_layout(i):
	picfile = "test.jpeg"
	prs = Presentation("default.pptx")
	graph_slide_layout = prs.slide_layouts[i]
	slide = prs.slides.add_slide(graph_slide_layout)
	placeholder = slide.placeholders[1]
	pic = placeholder.insert_picture(picfile)
	prs.save("../ppt/%s.pptx" % i)

def test_empty_layout():
	picfile = "test.jpeg"
	prs = Presentation("default.pptx")
	graph_slide_layout = prs.slide_layouts[6]
	slide = prs.slides.add_slide(graph_slide_layout)
	slide.shapes.add_picture(picfile,0,0)
	prs.save("../ppt/%s.pptx" % 6)
if __name__=="__main__":
	#test_empty_layout()
	#test_layout(int(sys.argv[1]))
	#test_ppt()

	crawler = Crawler()
	#crawler.get_item_list()
	item_url = "http://mp.weixin.qq.com/s?__biz=MzA3NzIwMDM3OA==&amp;mid=206906414&amp;idx=1&amp;sn=484555cf9c8efd164d06f6f6d0a6c19e&amp;scene=21#wechat_redirect"
	item_url = item_url.replace("&amp;","&")
	crawler.get_item(item_url)
	print "done!"